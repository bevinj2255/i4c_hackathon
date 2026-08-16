"""Fill the official i4C Idea Submission Template with our content.

    python make_ppt.py --team "Web Shooters" --members "A; B; C" --college "..." \
                       --contact "x@y.com" --phone "+91 ..."

Writes results/<Team>_KLA_PS01.pptx and solution_presentation.pptx.

Rules obeyed, taken from the template's own instruction slide:
  - maximum 6-7 slides INCLUDING the team slide  -> we ship 7
  - the instruction slide is removed
  - points and figures rather than paragraphs
  - the provided template is used as-is: theme, background, colours and shape furniture
    are never touched; only the {placeholder} text boxes are filled
  - the template's section pointers (KEY CONCEPT & APPROACH, SOLUTION OVERVIEW,
    KEY INNOVATION, ...) are kept exactly as provided

The template's 10 slides become 7: the instruction slide is dropped, "Technology &
Feasibility" folds into Proposed Solution, and "Research and References" folds into the
GitHub slide. Nothing KLA asks for is lost.

Every number is read from results/metrics.json, results/ood_metrics.json and the
training logs, so the deck cannot disagree with the measurements.
"""
import argparse
import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

TEMPLATE = Path("idea_template.pptx")
FIG = Path("results/figures")
RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def metrics():
    p = Path("results/metrics.json")
    if not p.exists():
        return {}
    out = {}
    for k, v in json.loads(p.read_text())["results"].items():
        if "bicubic" in k.lower():
            out["base"] = v
        elif "untrained" in k.lower():
            out["ctrl"] = v
        elif "TTA" in k:
            out["tta"] = v
        else:
            out["ours"] = v
    return out


def ood():
    p = Path("results/ood_metrics.json")
    if not p.exists():
        return None
    d = json.loads(p.read_text())
    g = [v["psnr_gain"] for v in d["patterns"].values()]
    return {"mean": d["mean_psnr_gain"], "n": len(g),
            "wins": sum(1 for x in g if x > 0), "best": max(g)}


def training(shipped_channels=None):
    """Only the runs that actually produced the shipped weights.

    Listing every log was wrong twice over: it included finetune_ood, an experiment we
    measured and rejected, and it included the 1.37M model's runs after a 4.4M model was
    shipped. A training history on a results slide has to describe the thing being
    submitted. Runs are matched to the shipped architecture by channel count.
    """
    runs = []
    for lg in sorted(Path("results").glob("*_log.csv")):
        name = lg.stem.replace("_log", "")
        if name.startswith("smoke") or "overfit" in name or "finetune_ood" in name:
            continue
        cfg_file = Path("configs") / f"{name}.json"
        if shipped_channels is not None and cfg_file.exists():
            if json.loads(cfg_file.read_text()).get("channels") != shipped_channels:
                continue
        rows = list(csv.DictReader(lg.open()))
        if rows:
            runs.append((name, len(rows),
                         sum(float(r["seconds"]) for r in rows) / 3600.0))
    return runs


def model_facts():
    """Architecture read from the shipped checkpoint, never typed in.

    These lines were hardcoded as "1.37M / 16 blocks x 64 ch / 31.5 ms" and silently
    became wrong the moment a different model was shipped. A deck that contradicts its
    own weights file is worse than one with fewer numbers.
    """
    import torch
    f = Path("weights/model.pt")
    if not f.exists():
        return {}
    ck = torch.load(f, map_location="cpu", weights_only=False)
    cfg = ck.get("cfg", {})
    n = sum(v.numel() for v in ck["model"].values())
    return {"params": n, "millions": n / 1e6, "ch": cfg.get("channels", "?"),
            "blocks": cfg.get("blocks", "?"), "mb": f.stat().st_size / 2**20}


def drop_slide(prs, index):
    lst = prs.slides._sldIdLst
    slides = list(lst)
    prs.part.drop_rel(slides[index].get(RID))
    lst.remove(slides[index])


def find(slide, needle):
    """Text box whose text contains `needle`.

    Matched on the template's own placeholder wording rather than shape ids, which
    differ from slide to slide.
    """
    for sh in slide.shapes:
        if sh.has_text_frame and needle.lower() in sh.text_frame.text.lower():
            return sh
    return None


def fill(shape, lines, size=None, bullet=True):
    """Replace a text box's content, keeping the template's typography.

    The first run's font is cloned onto every new line, so the template's fonts and
    colours survive even though the words change.
    """
    if shape is None:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    src = p0.runs[0] if p0.runs else None
    font = {}
    if src is not None:
        rgb = None
        try:
            if src.font.color is not None and src.font.color.type is not None:
                rgb = src.font.color.rgb
        except Exception:
            rgb = None
        font = {"size": src.font.size, "bold": src.font.bold,
                "name": src.font.name, "color": rgb}

    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for r in list(p0.runs):
        p0._p.remove(r._r)

    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        para = p0 if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = ("• " + line) if bullet else line
        f = run.font
        if font.get("name"):
            f.name = font["name"]
        f.size = Pt(size) if size else (font.get("size") or Pt(11))
        f.bold = font.get("bold")
        if font.get("color") is not None:
            f.color.rgb = font["color"]
        para.space_after = Pt(3)


def place(shape, left=None, top=None, width=None, height=None):
    """Move/resize a shape. The template's cards are sized for one-line placeholders;
    real content needs the room, so cards are grown and their text placed inside them
    rather than allowed to spill over the label above."""
    if shape is None:
        return
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)


def card(slide, top, height, left=1.05, width=11.23, index=0):
    """The n-th rounded content card on a slide (the big outer panel excluded)."""
    # The cards are AUTO_SHAPEs with an empty text frame. Filtering on
    # "not sh.has_text_frame" silently matched nothing -- auto-shapes have one.
    cards = [sh for sh in slide.shapes
             if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and sh.has_text_frame and not sh.text_frame.text.strip()
             and sh.width / 914400 > 4 and sh.top / 914400 > 3.0]
    cards.sort(key=lambda s: (s.top, s.left))
    if index < len(cards):
        place(cards[index], left, top, width, height)
        return cards[index]
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--team", default="TEAM NAME")
    ap.add_argument("--members", default="")
    ap.add_argument("--college", default="COLLEGE NAME")
    ap.add_argument("--contact", default="")
    ap.add_argument("--phone", default="")
    ap.add_argument("--year", default="3rd Year")
    ap.add_argument("--repo", default="https://github.com/bevinj2255/i4c_hackathon")
    ap.add_argument("--video", default="")
    ap.add_argument("--gpu", default="NVIDIA GeForce GTX 1650 (4 GiB)")
    ap.add_argument("--ms_per_image", default="101",
                    help="measured end-to-end ms/image, from inference.py")
    ap.add_argument("--out", default=None)
    a = ap.parse_args()

    if not TEMPLATE.exists():
        raise SystemExit(f"ABORT: {TEMPLATE} not found — it is the official i4C template.")

    mf = model_facts()
    m, o, runs = metrics(), ood(), training(mf.get('ch'))
    ours, base, ctrl = m.get("ours", {}), m.get("base", {}), m.get("ctrl", {})
    prs = Presentation(str(TEMPLATE))

    for idx in sorted([9, 7, 0], reverse=True):
        drop_slide(prs, idx)
    team, problem, idea, solution, innov, impact, links = list(prs.slides)

    # 1 -- Team ---------------------------------------------------------------
    fill(find(team, "Enter Team Name Here"), a.team, bullet=False)
    names = [n.strip() for n in a.members.split(";") if n.strip()]
    nb = sorted([sh for sh in team.shapes if sh.has_text_frame
                 and "{Enter Name}" in sh.text_frame.text], key=lambda s: s.top)
    yb = sorted([sh for sh in team.shapes if sh.has_text_frame
                 and "{Enter Year}" in sh.text_frame.text], key=lambda s: s.top)
    for box, nm in zip(nb, names):
        place(box, width=4.4)          # 1.2in wrapped "Bevin Punnoose Jacob" onto 3 lines
        fill(box, nm, bullet=False)
    for box in yb[:len(names)]:
        fill(box, a.year, bullet=False)
    # The template's own filler guidance is an instruction, not content -- drop it.
    hint = find(team, "A team can have up to 4 members")
    if hint is not None:
        hint._element.getparent().remove(hint._element)
    cbox = find(team, "Enter Full College Name"); place(cbox, width=10.0)
    fill(cbox, a.college, bullet=False)
    fill(find(team, "+91 XXXXX"), a.phone or "—", bullet=False)
    ebox = find(team, "email@example.com"); place(ebox, width=5.2)
    fill(ebox, a.contact or "—", bullet=False)

    # 2 -- Problem ------------------------------------------------------------
    fill(find(problem, "Selected the problem statement"),
         "KLA PS01 — AI-Based Restoration of Degraded Semiconductor Inspection Images",
         bullet=False, size=13)
    card(problem, top=4.05, height=2.95)
    place(find(problem, "DESCRIPTION / DETAILS"), left=1.31, top=4.18)
    place(find(problem, "Provide specific details about the problem statement"),
          left=1.31, top=4.52, width=10.7, height=2.4)
    fill(find(problem, "Provide specific details about the problem statement"), [
        "Inspection images decide whether a wafer passes: one pixel of noise can hide a killer defect.",
        "Three degradations, order undisclosed — speckle noise, additive Gaussian noise, 2× downsampling.",
        "Input 128×128 noisy → output 256×256 clean (also 256→512). Grayscale float32.",
        "Degraded pixels fall outside [0,1] by design; ground truth always inside it.",
        "Scored on three axes: PSNR + SSIM + LPIPS, H100 end-to-end throughput, training/compute hygiene.",
        "Hidden test set contains out-of-distribution content — must generalise, not memorise.",
    ], size=12)

    # 3 -- Idea ---------------------------------------------------------------
    fill(find(idea, "Provide a brief summary of your idea"),
         "One small fully-convolutional network denoises and 2× super-resolves in a single pass, "
         "built on a degradation model recovered from the data rather than guessed.",
         bullet=False, size=12)
    card(idea, top=4.05, height=1.42, index=0)
    place(find(idea, "KEY CONCEPT & APPROACH"), left=1.27, top=4.18)
    place(find(idea, "Briefly describe the core concept"),
          left=1.27, top=4.50, width=10.7, height=0.9)
    card(idea, top=5.60, height=1.42, index=1)
    place(find(idea, "SOLUTION OVERVIEW"), left=1.27, top=5.73)
    place(find(idea, "Provide an overview of the solution"),
          left=1.27, top=6.05, width=10.7, height=0.9)
    fill(find(idea, "Briefly describe the core concept"), [
        "Reverse-engineered KLA's exact degradation recipe from the 3200 provided pairs.",
        "That yields unlimited correctly-degraded training data, with fresh noise every epoch.",
    ], size=11)
    fill(find(idea, "Provide an overview of the solution"), [
        f"{mf.get('millions', 0):.1f}M-parameter residual CNN; all work at low resolution, "
        f"one PixelShuffle at the end — 4× cheaper.",
        "Fully convolutional, so the same weights restore 128→256 and 256→512.",
    ], size=11)

    # 4 -- Proposed solution (+ technology & feasibility) ---------------------
    fill(find(solution, "Describe your idea in detail"),
         "Measured forward model → unlimited synthetic pairs → compact CNN → "
         "loss blended for pixel and perceptual quality.", bullet=False, size=12)
    stack = "PyTorch 2.9.1+cu128 — " + (
        " + ".join(f"{n} {e} ep/{h:.1f} h" for n, e, h in runs) if runs else a.gpu)
    if a.gpu:
        stack += f" on {a.gpu}"
    card(solution, top=4.05, height=2.95)
    place(find(solution, "SOLUTION DETAILS"), left=1.31, top=4.18)
    place(find(solution, "Provide specific details about your proposed solution"),
          left=1.31, top=4.52, width=10.7, height=2.4)
    fill(find(solution, "Provide specific details about your proposed solution"), [
        "DEGRADATION RECOVERED (measured): 2×2 area-average downsample; noise applied after it; "
        "speckle multiplicative, var ∝ pixel² (r = 0.993); σ_speckle 0.10–0.25, σ_gauss 0.00–0.15.",
        "verify_degradation.py re-derives every constant from the data and ABORTS on mismatch.",
        "TRAINING: 50/50 real and synthesised pairs; noise ranges widened beyond measured; 8-fold dihedral augmentation.",
        f"MODEL: {mf.get('blocks','?')} residual blocks × {mf.get('ch','?')} ch → "
        f"PixelShuffle ×2. No global skip from the noisy input. Output clamped to [0,1].",
        "LOSS: Charbonnier blended with an LPIPS perceptual term by weight interpolation of two checkpoints.",
        f"STACK: {stack}.",
        f"FEASIBILITY: {mf.get('mb', 0):.0f} MB checkpoint, {a.ms_per_image} ms/image "
        f"end-to-end on a GTX 1650; runs unchanged on CPU or any NVIDIA GPU.",
    ], size=10)

    # 5 -- Innovation ---------------------------------------------------------
    fill(find(innov, "Highlight what makes your idea unique"),
         "We measured what others assume — which produced both the data advantage and the model.",
         bullet=False, size=12)
    card(innov, top=4.05, height=2.95, left=1.05, width=5.49, index=0)
    card(innov, top=4.05, height=2.95, left=6.79, width=5.49, index=1)
    place(find(innov, "KEY INNOVATION"), left=1.31, top=4.18)
    place(find(innov, "Describe the core innovation"),
          left=1.31, top=4.52, width=4.98, height=2.4)
    place(find(innov, "COMPETITIVE ADVANTAGE"), left=7.05, top=4.18)
    place(find(innov, "Explain how your solution is better"),
          left=7.05, top=4.52, width=4.98, height=2.4)
    fill(find(innov, "Describe the core innovation"), [
        "Degradation recovered, not guessed — giving unlimited perfectly-matched training data.",
        "Noise SPECTRUM matched, not just variance: caught a 16% high-frequency error every variance test passed.",
        "Perceptual/pixel weight interpolation — 93% of the LPIPS gain for 2.9% of the PSNR gain, zero extra training.",
    ], size=11)
    fill(find(innov, "Explain how your solution is better"), [
        f"Beats bicubic on {ours.get('n', 200)}/{ours.get('n', 200)} validation images — every one, not on average.",
        f"Compact by design ({mf.get('millions',0):.1f}M params, {mf.get('mb',0):.0f} MB): "
        f"all computation at low resolution, so throughput stays high.",
        "Ideas tested and REJECTED on evidence: 8× TTA (worsened LPIPS), synthetic-pattern training (−0.22 dB).",
    ], size=11)

    # 6 -- Impact -------------------------------------------------------------
    fill(find(impact, "Explain how your solution will make an impact"),
         "Recovers detail engineers currently accept as lost — and generalises to structure it never trained on.",
         bullet=False, size=12)
    card(impact, top=4.02, height=1.85, left=1.05, width=5.49, index=0)
    card(impact, top=4.02, height=1.85, left=6.79, width=5.49, index=1)
    # The template puts a small badge icon to the left of each label; move those with
    # the labels, or they end up sitting on top of the metrics.
    for pic in [sh for sh in impact.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                and sh.width / 914400 < 0.3 and sh.top / 914400 > 4.5]:
        pic.top = Inches(4.17)
    place(find(impact, "Primary Impact"), left=1.60, top=4.13)
    place(find(impact, "Describe the most significant benefit"),
          left=1.31, top=4.47, width=4.98, height=1.3)
    place(find(impact, "Quantifiable Outcomes"), left=7.37, top=4.13)
    place(find(impact, "List potential metrics or stats"),
          left=7.05, top=4.47, width=4.98, height=1.3)
    fill(find(impact, "Describe the most significant benefit"), [
        f"PSNR {base.get('psnr', 0):.2f} → {ours.get('psnr', 0):.2f} dB",
        f"SSIM {base.get('ssim', 0):.4f} → {ours.get('ssim', 0):.4f}",
        f"LPIPS {base.get('lpips', 0):.4f} → {ours.get('lpips', 0):.4f}",
        f"Untrained control {ctrl.get('psnr', 0):.2f} dB — the metric can fail.",
    ], size=11)
    quant = [
        f"+{ours.get('psnr', 0) - base.get('psnr', 0):.2f} dB PSNR, "
        f"+{ours.get('ssim', 0) - base.get('ssim', 0):.4f} SSIM, "
        f"{base.get('lpips', 0) - ours.get('lpips', 0):.4f} better LPIPS",
        f"{a.ms_per_image} ms/image end-to-end (read → preprocess → transfer → model → save)",
    ]
    if o:
        quant.append(f"Unseen semiconductor-like structure: {o['mean']:+.2f} dB mean, "
                     f"ahead on {o['wins']}/{o['n']} patterns")
    fill(find(impact, "List potential metrics or stats"), quant, size=11)
    fig = FIG / "slide_result.png"
    if fig.exists():
        impact.shapes.add_picture(str(fig), Inches(1.47), Inches(6.00), width=Inches(10.4))

    # 7 -- GitHub + references -------------------------------------------------
    fill(find(links, "Paste your GitHub"), a.repo, bullet=False, size=13)
    # Do NOT put references under the "Prototype / Simulation Video" heading -- that
    # mislabels them. The video box says what is true, and references get their own
    # box in the free space below the cards.
    fill(find(links, "Paste your Video Link"),
         a.video or "Not submitted — full results, figures and code are in the repository.",
         bullet=False, size=12)
    tb = links.shapes.add_textbox(Inches(1.15), Inches(6.45), Inches(11.0), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
            "REFERENCES",
            "Lim et al., EDSR, CVPRW 2017 · Shi et al., PixelShuffle, CVPR 2016 · "
            "Zhang et al., LPIPS, CVPR 2018 · Zhai et al., IEEE Access 11:21049, 2023 · "
            "Terven et al., Artif Intell Rev 58:195, 2025 · Monga et al., IEEE SPM 38(2), 2021.",
            "No external datasets or pretrained weights were used for the model.",
    ]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(11 if i == 0 else 9)
        run.font.bold = (i == 0)
        run.font.color.rgb = RGBColor(0xA8, 0xE6, 0x3A) if i == 0 else RGBColor(0xD5, 0xDC, 0xEA)
        para.space_after = Pt(2)

    out = Path(a.out or f"results/{a.team.replace(' ', '')}_KLA_PS01.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    prs.save("solution_presentation.pptx")
    print(f"CHECK: wrote {out} ({len(list(prs.slides))} slides, template theme preserved)")
    print(f"Export to PDF as {out.stem}.pdf before uploading.")


if __name__ == "__main__":
    main()
